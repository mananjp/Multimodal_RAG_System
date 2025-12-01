# -*- coding: utf-8 -*-
"""
Streamlit App: Multimodal RAG using CLIP + FAISS + Groq

Usage:
  1. Create .env file with: GROQ_API_KEY=gsk_your_key_here
  2. Run: streamlit run app.py

Features:
  - Upload multiple file formats (PDF, DOCX, PPTX, XLSX, TXT, Code, etc.)
  - Extract text and images using CLIP embeddings
  - Query documents with multimodal RAG
  - Display retrieved context and LLM responses
  - Cross-platform file handling (Windows, Mac, Linux)
  - Multiple LLM model selection per query
  - Universal file format support
"""

import os
import io
import tempfile
import streamlit as st
from pathlib import Path
from dotenv import load_dotenv
from langchain_groq import ChatGroq
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.embeddings import Embeddings
from sentence_transformers import SentenceTransformer
import json

# Load environment variables from .env file FIRST
load_dotenv()

# Verify API key is configured
if not os.environ.get("GROQ_API_KEY"):
    st.error(
        "❌ GROQ_API_KEY not found! "
        "Please create a .env file with: GROQ_API_KEY=gsk_your_key_here"
    )
    st.stop()

# Import RAG functions (assumed to be in multimodal_groq.py)
from multimodal_groq import (
    init_clip_model,
    process_pdf,
    build_vector_store,
    retrieve_multimodal,
    create_multimodal_message,
)

# Optional imports for different file formats
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

# -------------------------------------------------------------------
# Model Configuration
# -------------------------------------------------------------------

GROQ_MODELS = {
    "openai/gpt-oss-120b": "🚀 GPT-OSS 120B (Most Powerful)",
    "llama-3.1-8b-instant": "⚡ Llama 3.1 8B (Fast, Default)",
    "qwen/qwen3-32b": "🧠 Qwen3 32B (Balanced)",
    "playai-tts": "🎙️ PlayAI TTS (Audio)"
}

SUPPORTED_FORMATS = {
    "pdf": "📄 PDF Files",
    "docx": "📝 Word Documents",
    "pptx": "🎬 PowerPoint Presentations",
    "xlsx": "📊 Excel Spreadsheets",
    "txt": "📋 Text Files",
    "md": "📄 Markdown Files",
    "py": "🐍 Python Code",
    "java": "☕ Java Code",
    "js": "📜 JavaScript Code",
    "ts": "📘 TypeScript Code",
    "cpp": "⚙️ C++ Code",
    "c": "⚙️ C Code",
    "go": "🐹 Go Code",
    "rs": "🦀 Rust Code",
    "rb": "💎 Ruby Code",
    "php": "🐘 PHP Code",
    "cs": "💙 C# Code",
    "json": "🔧 JSON Files",
    "xml": "📋 XML Files",
    "html": "🌐 HTML Files",
    "yaml": "⚙️ YAML Files",
    "yml": "⚙️ YAML Files",
}

text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)

# -------------------------------------------------------------------
# Embeddings & Vector Store
# -------------------------------------------------------------------

class SentenceTransformerEmbeddings(Embeddings):
    """Wrapper for sentence-transformers to work with LangChain's FAISS."""

    def __init__(self, model_name: str = "all-MiniLM-L6-v2"):
        self.model = SentenceTransformer(model_name)

    def embed_documents(self, texts):
        return self.model.encode(texts, convert_to_numpy=True).tolist()

    def embed_query(self, text):
        return self.model.encode(text, convert_to_numpy=True).tolist()


@st.cache_resource
def load_text_embeddings():
    """Load sentence-transformers embeddings for text-only files."""
    return SentenceTransformerEmbeddings("all-MiniLM-L6-v2")


def build_text_vector_store(all_docs, embeddings_model):
    """Build FAISS vector store for text documents using sentence-transformers."""
    try:
        vector_store = FAISS.from_documents(all_docs, embeddings_model)
        return vector_store
    except Exception as e:
        st.error(f"Error building vector store: {str(e)}")
        return None


def retrieve_text_only(query, vector_store, k=5):
    """Retrieve documents using text-only embeddings (for non-PDF files)."""
    try:
        docs = vector_store.similarity_search(query, k=k)
        return docs
    except Exception as e:
        st.error(f"Error retrieving documents: {str(e)}")
        return []

# -------------------------------------------------------------------
# Message Builders
# -------------------------------------------------------------------

def create_text_message(query, context_docs):
    """Create a text-only message for non-PDF files."""
    context_text = "\n\n".join(
        f"Source: {doc.metadata.get('source', 'unknown')}\n{doc.page_content}"
        for doc in context_docs
    )

    message_content = f"""You are a helpful assistant analyzing documents.

Based on the following context from the document:

{context_text}

Please answer this question: {query}

Provide a clear, concise answer based on the context provided."""

    return HumanMessage(content=message_content)

# -------------------------------------------------------------------
# File Format Processors
# -------------------------------------------------------------------

def process_docx(file_path):
    """Process DOCX file and extract text."""
    if DocxDocument is None:
        raise ImportError("python-docx not installed. Run: pip install python-docx")

    doc = DocxDocument(file_path)
    all_docs = []

    for para_idx, para in enumerate(doc.paragraphs):
        if para.text.strip():
            doc_obj = Document(
                page_content=para.text,
                metadata={"page": para_idx, "type": "text", "source": "docx"}
            )
            all_docs.append(doc_obj)

    return all_docs


def process_pptx(file_path):
    """Process PPTX file and extract text from slides."""
    if Presentation is None:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(file_path)
    all_docs = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                doc_obj = Document(
                    page_content=shape.text,
                    metadata={"page": slide_idx, "type": "text", "source": "pptx"}
                )
                all_docs.append(doc_obj)

    return all_docs


def process_xlsx(file_path):
    """Process XLSX file and extract text from cells."""
    if openpyxl is None:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")

    wb = openpyxl.load_workbook(file_path)
    all_docs = []

    for sheet_idx, sheet in enumerate(wb.sheetnames):
        ws = wb[sheet]
        for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip() and row_text.strip() != "|":
                doc_obj = Document(
                    page_content=row_text,
                    metadata={"page": row_idx, "sheet": sheet, "type": "text", "source": "xlsx"}
                )
                all_docs.append(doc_obj)

    return all_docs


def process_text(file_path):
    """Process plain text, markdown, JSON, XML, HTML, YAML files."""
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    chunks = text_splitter.split_text(content)
    all_docs = []

    for idx, chunk in enumerate(chunks):
        doc_obj = Document(
            page_content=chunk,
            metadata={"page": idx, "type": "text", "source": "text"}
        )
        all_docs.append(doc_obj)

    return all_docs


def process_code(file_path):
    """Process code files (Python, Java, JavaScript, etc.)."""
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    chunks = text_splitter.split_text(content)
    all_docs = []

    for idx, chunk in enumerate(chunks):
        doc_obj = Document(
            page_content=chunk,
            metadata={"page": idx, "type": "text", "source": "code"}
        )
        all_docs.append(doc_obj)

    return all_docs


def get_file_extension(filename):
    """Get file extension from filename."""
    return Path(filename).suffix.lstrip(".").lower()


def process_file(file_path, filename):
    """
    Universal file processor - detects format and routes to appropriate handler.
    Returns: all_docs for consistent interface
    """
    ext = get_file_extension(filename)

    if ext == "docx":
        all_docs = process_docx(file_path)
    elif ext == "pptx":
        all_docs = process_pptx(file_path)
    elif ext == "xlsx":
        all_docs = process_xlsx(file_path)
    elif ext in ["txt", "md", "json", "xml", "html", "yaml", "yml"]:
        all_docs = process_text(file_path)
    elif ext in ["py", "java", "js", "ts", "cpp", "c", "go", "rs", "rb", "php", "cs"]:
        all_docs = process_code(file_path)
    else:
        raise ValueError(f"Unsupported file format: .{ext}")

    return all_docs

# -------------------------------------------------------------------
# Streamlit Page Configuration
# -------------------------------------------------------------------

st.set_page_config(
    page_title="Multimodal RAG",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown(
    """
    <style>
        .main { max-width: 1200px; margin: 0 auto; }
        .stChatMessage { border-radius: 8px; padding: 12px; }
        .retrieved-doc { background-color: #f0f2f6; padding: 10px; border-radius: 5px; margin: 5px 0; }
        .doc-text { color: #0f1419; font-size: 14px; }
        .doc-image { color: #ff6b6b; font-weight: bold; }
        .model-badge { display: inline-block; padding: 4px 8px; border-radius: 4px; font-size: 12px; font-weight: bold; margin-left: 8px; }
        .file-type-badge { display: inline-block; padding: 2px 6px; border-radius: 3px; font-size: 10px; margin-right: 5px; background: #e3f2fd; color: #1976d2; }
    </style>
    """,
    unsafe_allow_html=True,
)

# -------------------------------------------------------------------
# Session State Initialization
# -------------------------------------------------------------------

@st.cache_resource
def load_clip_model():
    """Cache CLIP model to avoid reloading."""
    with st.spinner("⏳ Loading CLIP model (first time only)..."):
        return init_clip_model()


def get_llm_for_model(model_name):
    """Get or create cached LLM instance for a specific model."""
    if "llm_instances" not in st.session_state:
        st.session_state.llm_instances = {}
    if model_name not in st.session_state.llm_instances:
        st.session_state.llm_instances[model_name] = ChatGroq(
            model_name=model_name,
            temperature=0.7,
        )
    return st.session_state.llm_instances[model_name]


def init_session_state():
    """Initialize Streamlit session state."""
    if "vector_store" not in st.session_state:
        st.session_state.vector_store = None
    if "file_processed" not in st.session_state:
        st.session_state.file_processed = False
    if "file_name" not in st.session_state:
        st.session_state.file_name = None
    if "file_type" not in st.session_state:
        st.session_state.file_type = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "last_query" not in st.session_state:
        st.session_state.last_query = None
    if "selected_model" not in st.session_state:
        st.session_state.selected_model = "llama-3.1-8b-instant"
    if "llm_instances" not in st.session_state:
        st.session_state.llm_instances = {}
    if "all_docs" not in st.session_state:
        st.session_state.all_docs = []
    if "clip_model" not in st.session_state:
        st.session_state.clip_model = None
    if "clip_processor" not in st.session_state:
        st.session_state.clip_processor = None


init_session_state()

# -------------------------------------------------------------------
# Sidebar Configuration
# -------------------------------------------------------------------

st.sidebar.markdown("## ⚙️ Configuration")

with st.sidebar.expander("API & Model Settings", expanded=True):
    api_key = os.environ.get("GROQ_API_KEY", "")
    if api_key:
        masked_key = api_key[:10] + "..." + api_key[-4:]
        st.success(f"✅ API Key Configured: {masked_key}")
        st.info("💡 All requests routed through preconfigured API key.")
    else:
        st.error("❌ No API key found!")

    st.markdown("### 🤖 Default LLM Model")
    default_model = st.selectbox(
        "Select default model for queries:",
        options=list(GROQ_MODELS.keys()),
        format_func=lambda x: GROQ_MODELS[x],
        key="default_model_select",
    )
    st.session_state.selected_model = default_model

st.sidebar.markdown("---")
st.sidebar.markdown("## 📁 File Upload")

allowed_types = list(SUPPORTED_FORMATS.keys())
uploaded_file = st.sidebar.file_uploader(
    "Upload a file:",
    type=allowed_types,
    help="Supported: PDF, DOCX, PPTX, XLSX, TXT, Code files, Markdown, JSON, XML, HTML, YAML",
)

# -------------------------------------------------------------------
# File Processing
# -------------------------------------------------------------------

if uploaded_file is not None:
    file_bytes = uploaded_file.read()
    file_name = uploaded_file.name
    file_ext = get_file_extension(file_name)

    if file_name != st.session_state.file_name or not st.session_state.file_processed:
        st.sidebar.info(f"🔄 Processing {file_ext.upper()} file...")

        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                temp_file_path = os.path.join(tmpdir, file_name)
                with open(temp_file_path, "wb") as f:
                    f.write(file_bytes)

                if file_ext == "pdf":
                    clip_model, clip_processor = load_clip_model()
                else:
                    clip_model, clip_processor = None, None

                with st.spinner(f"📄 Processing {file_ext.upper()} file..."):
                    if file_ext == "pdf":
                        all_docs, embeddings_array, image_data_store = process_pdf(
                            temp_file_path, clip_model, clip_processor
                        )
                        vector_store = build_vector_store(all_docs, embeddings_array)
                    else:
                        all_docs = process_file(temp_file_path, file_name)
                        embeddings = load_text_embeddings()
                        vector_store = build_text_vector_store(all_docs, embeddings)

                st.session_state.vector_store = vector_store
                st.session_state.clip_model = clip_model
                st.session_state.clip_processor = clip_processor
                st.session_state.file_processed = True
                st.session_state.file_name = file_name
                st.session_state.file_type = file_ext
                st.session_state.all_docs = all_docs

                st.sidebar.success(f"✅ {file_ext.upper()} processed! ({len(all_docs)} chunks)")
                st.sidebar.info(f"📄 File: {file_name}")

        except Exception as e:
            st.sidebar.error(f"❌ Error processing file: {str(e)}")
            import traceback
            st.sidebar.error(f"Details: {traceback.format_exc()}")

# -------------------------------------------------------------------
# Main Chat Interface
# -------------------------------------------------------------------

st.title("📚 Multimodal RAG - Universal Document Q&A")

if not st.session_state.file_processed:
    st.info(
        "👈 Please upload a file in the sidebar to get started. "
        "Supports: PDF, DOCX, PPTX, XLSX, TXT, Code files, Markdown, JSON, XML, HTML, YAML",
        icon="ℹ️",
    )
else:
    file_type_emoji = SUPPORTED_FORMATS.get(st.session_state.file_type, "📄")
    st.success(f"✅ Ready! Loaded: **{st.session_state.file_name}** {file_type_emoji}")

    col1, col2, col3 = st.columns([2, 1, 1.2])

    with col1:
        query = st.text_input(
            "💬 Ask a question about the document:",
            placeholder="e.g., Summarize the main findings...",
            key="query_input",
        )

    with col2:
        num_chunks = st.number_input(
            "Retrieve chunks:",
            min_value=1,
            max_value=10,
            value=5,
            step=1,
        )

    with col3:
        selected_model = st.selectbox(
            "LLM Model:",
            options=list(GROQ_MODELS.keys()),
            format_func=lambda x: GROQ_MODELS[x],
            key="query_model_select",
            index=list(GROQ_MODELS.keys()).index(st.session_state.selected_model),
        )

    if query and st.session_state.file_processed and query != st.session_state.last_query:
        st.session_state.last_query = query

        with st.spinner("🔍 Searching document & generating response..."):
            try:
                if st.session_state.file_type == "pdf":
                    context_docs = retrieve_multimodal(
                        query,
                        st.session_state.vector_store,
                        st.session_state.clip_model,
                        st.session_state.clip_processor,
                        k=num_chunks,
                    )
                else:
                    context_docs = retrieve_text_only(
                        query,
                        st.session_state.vector_store,
                        k=num_chunks,
                    )

                llm = get_llm_for_model(selected_model)

                if st.session_state.file_type == "pdf":
                    message = create_multimodal_message(query, context_docs)
                else:
                    message = create_text_message(query, context_docs)

                response = llm.invoke([message])
                answer = response.content

                st.session_state.chat_history.append(
                    {
                        "query": query,
                        "answer": answer,
                        "docs": context_docs,
                        "model": selected_model,
                    }
                )

            except Exception as e:
                st.error(f"❌ Error processing query: {str(e)}")
                import traceback
                st.error(f"Details: {traceback.format_exc()}")
                answer = None

        if answer:
            st.markdown("---")
            st.markdown(f"### 🤖 Response {GROQ_MODELS[selected_model]}")
            st.markdown(answer)

            st.markdown("### 📋 Retrieved Context")

            text_docs = [d for d in context_docs if d.metadata.get("type") == "text"]
            image_docs = [d for d in context_docs if d.metadata.get("type") == "image"]

            if text_docs:
                st.markdown("#### Text Excerpts:")
                for i, doc in enumerate(text_docs, 1):
                    page = doc.metadata.get("page", "?")
                    with st.expander(
                        f"📖 Text from Page {page} (chunk {i})", expanded=(i == 1)
                    ):
                        st.markdown(
                            f'<div class="retrieved-doc"><p class="doc-text">{doc.page_content}</p></div>',
                            unsafe_allow_html=True,
                        )

            if image_docs:
                st.markdown("#### Images Found:")
                for doc in image_docs:
                    page = doc.metadata.get("page", "?")
                    st.markdown(
                        f'<div class="retrieved-doc"><p class="doc-image">🖼️ Image from Page {page}</p></div>',
                        unsafe_allow_html=True,
                    )

    if st.session_state.chat_history:
        st.markdown("---")
        st.markdown("### 📜 Chat History")

        for i, exchange in enumerate(st.session_state.chat_history[::-1]):
            col1, col2 = st.columns([10, 1])
            with col1:
                model_name = exchange.get("model", "unknown")
                model_label = GROQ_MODELS.get(model_name, model_name)
                st.markdown(
                    f"**Q {len(st.session_state.chat_history) - i}:** {exchange['query']}"
                    f"  <span class='model-badge'>{model_label}</span>",
                    unsafe_allow_html=True,
                )
                st.markdown(f"> {exchange['answer'][:200]}...")

            with col2:
                if st.button("↻", key=f"repeat_{i}"):
                    st.session_state.last_query = None

        if st.button("🗑️ Clear History"):
            st.session_state.chat_history = []
            st.session_state.last_query = None

# -------------------------------------------------------------------
# Footer
# -------------------------------------------------------------------

st.markdown("---")
st.markdown(
    """
    <div style="text-align: center; color: #666; font-size: 12px;">
    Multimodal RAG powered by CLIP + FAISS + Groq 🚀 | Multi-format support | Multi-model LLM
    </div>
    """,
    unsafe_allow_html=True,
)
